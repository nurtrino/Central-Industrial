"""
Deep Research — standalone local server.

Extracted from the original platform's perf_server.py (port 5002) so the
Deep Research Tool (DRT) lives on its own, as a "Special Projects" tool. Unlike the
old setup (GitHub-Pages HTTPS frontend → localhost backend, hence CORS), this server
SERVES ITS OWN UI at / and exposes the DRT API on the SAME origin — like the hub and
Monkey Read Monkey Do. Self-contained: imports engines.research.* (copied verbatim);
config/, prompts/, and the encrypted vault sit beside this file so the engine's
relative path math (dirname×3 → repo root) resolves unchanged.

PORT : 5006
Start: python dr_server.py   (or Deep Research.vbs / the hub's "launch" on demand)
"""

import base64
import hashlib
import hmac
import io
import json
import os
import re
import sys
import tempfile
import threading
import time
import traceback

from flask import Flask, jsonify, redirect, request, send_from_directory
from urllib.parse import urlencode

_ROOT = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _ROOT)
_PROMPTS_DIR = os.path.join(_ROOT, "prompts")

# PORT: Render injects $PORT at runtime; default 5006 for local runs.
PORT = int(os.environ.get("PORT", "5006"))


# ── Load .env at startup so the server works when launched via pythonw / the hub
# (no shell environment). The guard only sets a key that is MISSING or EMPTY —
# never overwrite a real value (Claude Code injects an empty key into subprocesses;
# this prevents that from clobbering the real ANTHROPIC_API_KEY). Do not revert this.
def _load_dotenv():
    env_path = os.path.join(_ROOT, ".env")
    if not os.path.exists(env_path):
        return
    with open(env_path, encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            key, _, val = line.partition("=")
            key = key.strip()
            val = val.strip().strip('"').strip("'")
            if key and not os.environ.get(key, "").strip():  # set if missing or empty
                os.environ[key] = val

_load_dotenv()

AUTH_SECRET = os.environ.get("AUTH_SECRET", "")
HOME_URL = (os.environ.get("HOME_URL", "").strip()
            or os.environ.get("HUB_URL", "http://127.0.0.1:5050/").strip())
GATE_ON = bool(AUTH_SECRET)

app = Flask(__name__)


# ── Headers: same-origin now, but keep CORS permissive (harmless) and force
# no-store so the served page is never stale — mirrors the hub.
@app.after_request
def _headers(resp):
    resp.headers["Access-Control-Allow-Origin"]  = "*"
    resp.headers["Access-Control-Allow-Methods"] = "GET, POST, OPTIONS"
    resp.headers["Access-Control-Allow-Headers"] = "Content-Type"
    resp.headers["Cache-Control"] = "no-store, must-revalidate"
    resp.headers["Expires"] = "0"
    return resp


# ── UI (served from disk on every request) ───────────────────────────────────
SESS_TTL = 8 * 3600


def _verify(purpose, tok):
    try:
        exp_s, sig = (tok or "").split(".", 1)
        exp = int(exp_s)
    except (ValueError, AttributeError):
        return False
    if exp < int(time.time()):
        return False
    good = hmac.new(AUTH_SECRET.encode(), f"{purpose}:{exp}".encode(), hashlib.sha256).hexdigest()
    return hmac.compare_digest(sig, good)


def _make_sess():
    exp = int(time.time()) + SESS_TTL
    sig = hmac.new(AUTH_SECRET.encode(), f"sess:{exp}".encode(), hashlib.sha256).hexdigest()
    return f"{exp}.{sig}"


def _authed():
    return _verify("sess", request.cookies.get("ci_sess", ""))


@app.before_request
def _access_gate():
    if not GATE_ON or request.method == "OPTIONS":
        return
    p = request.path
    if p == "/api/health":
        return
    sso = request.args.get("t", "")
    if sso and _verify("sso", sso):
        rest = request.args.to_dict(flat=True)
        rest.pop("t", None)
        clean = p + ("?" + urlencode(rest) if rest else "")
        resp = redirect(clean, code=302)
        resp.set_cookie("ci_sess", _make_sess(), max_age=None,
                        httponly=True, secure=True, samesite="Lax")
        return resp
    if _authed():
        return
    if p.startswith("/api/"):
        return jsonify({"error": "unauthorized"}), 401
    return redirect(HOME_URL, code=302)


# ── UI (served from disk on every request) ──────────────────────────────
_HUB_URL = os.environ.get("HUB_URL", "").strip() or HOME_URL
if _HUB_URL and not re.match(r"^https?://", _HUB_URL, re.I):
    _HUB_URL = "https://" + _HUB_URL


@app.route("/")
def index():
    with open(os.path.join(_ROOT, "index.html"), encoding="utf-8") as fh:
        html = fh.read().replace("__HUB_URL__", _HUB_URL)
    return app.response_class(html, mimetype="text/html")


@app.route("/vendor/<path:fn>")
def vendor(fn):
    return send_from_directory(os.path.join(_ROOT, "vendor"), fn)


@app.route("/fonts/<path:fn>")
def fonts(fn):
    return send_from_directory(os.path.join(_ROOT, "fonts"), fn)


@app.route("/favicon.ico")
def favicon():
    return ("", 204)


_EGG_MP3 = "Boonies Basement Tub (128kbit_AAC)-2.mp3"


@app.route("/egg.mp3")
def egg_mp3():
    """The 'Oh, Very Deep' easter egg — played by the UI when that depth chip is selected.
    The file lives beside the server locally only (not committed); elsewhere this 404s and
    the UI's play() fails silently."""
    path = os.path.join(_ROOT, _EGG_MP3)
    if not os.path.exists(path):
        return ("", 404)
    return send_from_directory(_ROOT, _EGG_MP3, mimetype="audio/mpeg")


@app.route("/api/health")
def health():
    try:
        from engines.research import brightdata as _bd, tavily_search as _tv, exa_search as _ex
        from engines.research import brave_search as _br
        providers = {"brightdata": _bd.status(), "tavily": _tv.is_enabled(), "exa": _ex.is_enabled(),
                     "brave_api": _br.is_enabled()}
    except Exception:
        providers = {}
    return jsonify({"ok": True, "tool": "deep-research", "build": "2026-09-03.sonnetlanes",
                    "sources_hash": _sources_hash(), "providers": providers})


@app.route("/api/firecrawl/<path:fcpath>", methods=["GET", "POST", "OPTIONS"])
def firecrawl_proxy(fcpath):
    """Proxy for the Firecrawl sub-tool: forwards to https://api.firecrawl.dev/<path>,
    injecting the key server-side (never exposed to the browser). Pass-through JSON + status."""
    if request.method == "OPTIONS":
        return "", 204
    key = os.environ.get("FIRECRAWL_API_KEY", "").strip()
    if not key:
        return jsonify({"error": "FIRECRAWL_API_KEY is not configured — add it to .env and restart."}), 400
    import requests
    url = "https://api.firecrawl.dev/" + fcpath.lstrip("/")
    headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    try:
        if request.method == "POST":
            r = requests.post(url, headers=headers, json=(request.get_json(silent=True) or {}), timeout=300)
        else:
            r = requests.get(url, headers=headers, params=request.args, timeout=300)
    except Exception as e:  # noqa: BLE001
        return jsonify({"error": f"could not reach Firecrawl: {e}"}), 502
    try:
        return jsonify(r.json()), r.status_code
    except Exception:
        return (r.text, r.status_code, {"Content-Type": r.headers.get("Content-Type", "text/plain")})


@app.route("/api/restart", methods=["POST", "OPTIONS"])
def restart_endpoint():
    """Restart the server (re-reads .env + config/drt_models.json). Spawns a detached
    helper that waits for this process to free the port, relaunches, then this exits."""
    if request.method == "OPTIONS":
        return "", 204
    import time as _time
    import subprocess
    helper = os.path.join(_ROOT, "restart_helper.py")
    DETACHED = 0x00000008 | 0x00000200  # DETACHED_PROCESS | CREATE_NEW_PROCESS_GROUP
    try:
        subprocess.Popen([sys.executable, helper], cwd=_ROOT,
                         creationflags=DETACHED, close_fds=True)
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500
    def _bye():
        _time.sleep(0.8)
        os._exit(0)
    threading.Thread(target=_bye, daemon=True).start()
    return jsonify({"ok": True, "message": "Server restarting…"})


# ── Save a generated report (.docx) to disk + open it in Word ────────────────
_REPORTS_DIR = os.environ.get("DRT_REPORTS_DIR", os.path.join(_ROOT, "reports"))
_FN_STOPWORDS = {"the", "a", "an", "of", "and", "or", "to", "in", "on", "for", "what",
                 "is", "are", "how", "who", "whom", "does", "do", "did", "about", "any",
                 "with", "into", "at", "by", "from", "vs", "re", "&"}


def _slug_from_query(q: str) -> str:
    """A short, human-readable keyword slug drawn from the research topic."""
    words = re.findall(r"[A-Za-z0-9]+", q or "")
    keep = [w for w in words if w.lower() not in _FN_STOPWORDS and len(w) > 1]
    keep = keep[:6] or words[:4]
    return (" ".join(keep).strip())[:70].strip() or "report"


def _safe_filename(name: str) -> str:
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "", name or "")
    name = re.sub(r"\s+", " ", name).strip().strip(".")
    return name[:140] or "report"


def _report_title(query: str, report_md: str = "", client=None) -> str:
    """A subject-appropriate document title for the report — one cheap model call
    ("4-9 words, Title Case, specific to what the report actually found"), through the
    wrapped client. Falls back to the keyword slug on any failure, so a filename always
    exists. Result is filename-safe."""
    title = ""
    try:
        if client is None:
            from engines.research.llm import make_client
            client = make_client("claude", os.environ.get("ANTHROPIC_API_KEY", "").strip())
        if client is not None:
            from engines.research.models import get_model
            body = (report_md or "")[:6000]
            sys_p = ("You title finished research reports for a document library. Reply with "
                     "ONLY the title: 4-9 words, Title Case, specific to the subject and what "
                     "the report actually established (name the entity/product/version/topic — "
                     "not the question's phrasing), no quotes, no trailing period, no colon, "
                     "no 'Report on' / 'Analysis of' filler.")
            r = client.messages.create(
                model=get_model("route"), max_tokens=60, system=sys_p,
                messages=[{"role": "user", "content":
                           f"RESEARCH QUESTION:\n{query}\n\nREPORT (opening):\n{body}"}])
            title = "".join(getattr(b, "text", "") for b in r.content
                            if getattr(b, "type", "") == "text").strip()
            title = title.splitlines()[0].strip().strip('"\'*#').rstrip(".:").strip()
            if len(title.split()) > 12 or len(title) < 4:
                title = ""
    except Exception:
        title = ""
    if not title:
        title = _slug_from_query(query).title()
    return _safe_filename(title)[:90] or "Report"


def _report_filename(title: str, label: str, when=None) -> str:
    """'<Title> — 2026-09-03 1347.docx' (label appended only for non-default tools,
    e.g. '<Title> (Odysseus) — …'). The date stamp keeps same-day re-runs distinct."""
    import time as _time
    stamp = _time.strftime("%Y-%m-%d %H%M", _time.localtime(when) if when else _time.localtime())
    suffix = "" if (label or "").strip().lower() in ("", "deep research") else f" ({label.strip()})"
    return _safe_filename(f"{title}{suffix} — {stamp}") + ".docx"


def _autosave_report(docx_b64: str, query: str, label: str, title: str = "") -> str:
    """Write a finished report .docx straight to the reports folder the moment a run
    completes (same titled filename as /api/save_report), so the report is on disk even
    if the Save button is never clicked. Does NOT open Word. Never raises — returns
    the saved path, or '' on failure (the in-browser copy + Save button still work)."""
    try:
        raw = base64.b64decode(docx_b64 or "")
        if not raw:
            return ""
        os.makedirs(_REPORTS_DIR, exist_ok=True)
        fname = _report_filename(title or _report_title(query), label)
        path = os.path.join(_REPORTS_DIR, fname)
        with open(path, "wb") as fh:
            fh.write(raw)
        return path
    except Exception:
        return ""


@app.route("/api/save_report", methods=["POST", "OPTIONS"])
def save_report():
    r"""Save a generated .docx to D:\______Documents\___Deep Research Reports as
    '<subject title> — <date stamp>.docx', then open it in Word.
    Body: {docx_b64, query, label, title?, report_md?, open?}. `title` (the run's
    generated title) is used when given; otherwise one is generated from the report.
    Returns {ok, path, filename, opened}. (open=false skips the Word launch.)"""
    if request.method == "OPTIONS":
        return "", 204
    data = request.get_json(silent=True) or {}
    b64 = data.get("docx_b64") or ""
    query = (data.get("query") or "").strip()
    label = (data.get("label") or "Deep Research").strip()
    title = _safe_filename((data.get("title") or "").strip())[:90]
    do_open = data.get("open", True)
    if not b64:
        return jsonify({"error": "no document to save"}), 400
    try:
        raw = base64.b64decode(b64)
    except Exception:
        return jsonify({"error": "document data was not valid base64"}), 400
    try:
        os.makedirs(_REPORTS_DIR, exist_ok=True)
    except Exception as e:  # noqa: BLE001
        return jsonify({"error": f"could not create the reports folder: {e}"}), 500
    if not title:
        title = _report_title(query, data.get("report_md") or "")
    fname = _report_filename(title, label)
    path = os.path.join(_REPORTS_DIR, fname)
    try:
        with open(path, "wb") as fh:
            fh.write(raw)
    except Exception as e:  # noqa: BLE001
        return jsonify({"error": f"could not save the file: {e}"}), 500
    opened = False
    if do_open and hasattr(os, "startfile"):   # Windows only — hosted/Linux reports
        try:                                   # opened=False and the UI downloads it.
            os.startfile(path)   # opens in the registered .docx app (Word)
            opened = True
        except Exception:
            opened = False
    return jsonify({"ok": True, "path": path, "filename": fname, "opened": opened})


# ── Async job registry (standalone — no longer shared with other tools) ──────
_JOBS = {}
_JOBS_LOCK = threading.Lock()


# ── Per-job event stream (the UI's live activity feed) ───────────────────────
# Each job dict carries "events": [[seq, kind, text], …] and "eseq": <latest seq>.
# Workers push raw engine log lines here; the status endpoints return the slice
# with seq > ?after=<n> so the UI can poll incrementally.
_EVENT_CAP = 3000          # ring buffer: drop oldest beyond this
_EVENT_TEXT_CAP = 500      # per-event text cap

_EVENT_KIND_MAP = {
    "search": "search", "open": "open", "forum": "forum",
    "extension": "ext", "login": "login", "note": "note", "plan": "plan",
    # High-granularity cascade: every tool call, every hit, every turn, every
    # extraction verdict, every synthesis/meta step flows to the activity panel.
    "act": "act", "hit": "hit", "turn": "turn", "extract": "extract", "synth": "synth",
    "docs": "docs",   # uploaded-document intake analysis
    "lane": "lane",   # parallel-lane planning / lead review / hot-trail dig
    "brightdata": "search", "tavily": "search", "exa": "search", "local-baseline": "search",
    "deepen": "synth", "stage2": "stage", "stage3": "stage", "stage4": "stage",
}


def _event_kind(msg):
    """Derive an event kind from a log line's leading "[xxx]" prefix."""
    m = re.match(r"\s*\[([a-z_\-]+)\]", str(msg or ""), re.I)
    if m:
        return _EVENT_KIND_MAP.get(m.group(1).lower(), "log")
    return "log"


def _push_event(job_id, kind, text):
    """Append one event to a job's stream. Thread-safe; silently no-ops if the
    job is gone (finished/popped). Never raises."""
    try:
        with _JOBS_LOCK:
            job = _JOBS.get(job_id)
            if job is None:
                return
            seq = job.get("eseq", 0) + 1
            job["eseq"] = seq
            evs = job.setdefault("events", [])
            evs.append([seq, kind, str(text or "")[:_EVENT_TEXT_CAP]])
            if len(evs) > _EVENT_CAP:
                del evs[:len(evs) - _EVENT_CAP]
    except Exception:
        pass


# ── Shared helpers carried from perf_server.py ───────────────────────────────
def _extract_file_text(path: str, filename: str) -> str:
    """Extract plain text from a supporting document, by extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"

    if ext == "pdf":
        import pdfplumber
        pages = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
        return "\n\n".join(pages)

    elif ext in ("pptx", "ppt"):
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)

    elif ext in ("docx", "doc"):
        import mammoth
        with open(path, "rb") as fh:
            result = mammoth.extract_raw_text(fh)
        return result.value

    elif ext in ("xlsx", "xls"):
        import pandas as pd
        sheets = []
        with pd.ExcelFile(path) as xf:
            for sname in xf.sheet_names:
                df = xf.parse(sname, header=None)
                sheets.append(f"[Sheet: {sname}]\n{df.to_string(index=False)}")
        return "\n\n".join(sheets)

    elif ext in ("csv",):
        import pandas as pd
        df = pd.read_csv(path)
        return df.to_string(index=False)

    else:
        with open(path, encoding="utf-8", errors="replace") as fh:
            return fh.read()


def _backup_then_write(path, content):
    """Write `content` to `path`, keeping a rolling single .bak of the prior version."""
    if os.path.exists(path):
        try:
            with open(path, encoding="utf-8") as _src, \
                 open(path + ".bak", "w", encoding="utf-8") as _dst:
                _dst.write(_src.read())
        except Exception:
            pass
    with open(path, "w", encoding="utf-8") as fh:
        fh.write(content)


def _memo_to_docx_bytes(memo_md: str, label: str, images: dict = None) -> bytes:
    """Convert markdown to a .docx (returned as bytes). Renders GitHub-flavored pipe
    tables as real Word tables, plus headings, bullets, rules, and inline **bold**/*italic*."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    for section in doc.sections:
        section.top_margin    = Inches(0.5)
        section.bottom_margin = Inches(0.5)
        section.left_margin   = Inches(0.5)
        section.right_margin  = Inches(0.5)

    HEADER_BG = "1A1A2E"
    _INLINE = re.compile(r"(\*\*.+?\*\*|__.+?__|\*.+?\*|_.+?_)")

    def _add_runs(paragraph, text: str):
        pos = 0
        for m in _INLINE.finditer(text):
            if m.start() > pos:
                paragraph.add_run(text[pos:m.start()])
            tok = m.group(0)
            if tok.startswith("**") or tok.startswith("__"):
                paragraph.add_run(tok[2:-2]).bold = True
            else:
                paragraph.add_run(tok[1:-1]).italic = True
            pos = m.end()
        if pos < len(text):
            paragraph.add_run(text[pos:])

    def _add_heading(text: str, level: int):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        run = p.add_run(text)
        run.bold = True
        run.font.size = Pt(13 if level == 1 else 11.5 if level == 2 else 11)
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        p.paragraph_format.space_before = Pt(12 if level == 1 else 8)
        p.paragraph_format.space_after  = Pt(2)

    def _add_body(text: str):
        p = doc.add_paragraph()
        _add_runs(p, text)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(4)

    def _add_bullet(text: str):
        p = doc.add_paragraph(style="List Bullet")
        _add_runs(p, text)
        p.paragraph_format.space_after = Pt(2)

    def _split_row(line: str):
        s = line.strip()
        if s.startswith("|"):
            s = s[1:]
        if s.endswith("|"):
            s = s[:-1]
        return [c.strip() for c in s.split("|")]

    def _is_sep_row(line: str) -> bool:
        cells = _split_row(line)
        if not cells:
            return False
        saw_dash = False
        for c in cells:
            cc = c.strip()
            if cc == "":
                continue
            if set(cc) <= set("-:") and "-" in cc:
                saw_dash = True
            else:
                return False
        return saw_dash

    def _set_cell_bg(cell, hex_color: str):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color)
        tcPr.append(shd)

    def _add_table(rows):
        ncols = max(len(r) for r in rows)
        tbl = doc.add_table(rows=0, cols=ncols)
        try:
            tbl.style = "Table Grid"
        except KeyError:
            pass
        tbl.autofit = True
        for ri, row in enumerate(rows):
            cells = tbl.add_row().cells
            for ci in range(ncols):
                val  = row[ci] if ci < len(row) else ""
                cell = cells[ci]
                cell.text = ""
                p = cell.paragraphs[0]
                p.paragraph_format.space_before = Pt(1)
                p.paragraph_format.space_after  = Pt(1)
                _add_runs(p, val)
                if ri == 0:
                    _set_cell_bg(cell, HEADER_BG)
                    for r in p.runs:
                        r.bold = True
                        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        r.font.size = Pt(10)
                else:
                    for r in p.runs:
                        r.font.size = Pt(10)
        spacer = doc.add_paragraph()
        spacer.paragraph_format.space_after = Pt(2)

    lines = memo_md.splitlines()
    n = len(lines)
    i = 0
    while i < n:
        stripped = lines[i].strip()
        if not stripped:
            i += 1
            continue
        if "|" in stripped and i + 1 < n and _is_sep_row(lines[i + 1]):
            rows = [_split_row(lines[i])]
            j = i + 2
            while j < n and "|" in lines[j] and lines[j].strip().startswith("|"):
                rows.append(_split_row(lines[j]))
                j += 1
            _add_table(rows)
            i = j
            continue
        if re.fullmatch(r"-{3,}|\*{3,}|_{3,}", stripped):
            i += 1
            continue
        m_img = re.fullmatch(r"\{\{IMG:(.+?)\}\}", stripped)
        if m_img and images and m_img.group(1) in images:
            try:
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.add_run().add_picture(
                    io.BytesIO(images[m_img.group(1)]["bytes"]), width=Inches(5.5)
                )
                p.paragraph_format.space_before = Pt(4)
                p.paragraph_format.space_after = Pt(4)
            except Exception:
                pass
            i += 1
            continue
        if stripped.startswith("### "):
            _add_heading(stripped[4:], 3)
        elif stripped.startswith("## "):
            _add_heading(stripped[3:], 2)
        elif stripped.startswith("# "):
            _add_heading(stripped[2:], 1)
        elif stripped.startswith(("- ", "* ", "• ")):
            _add_bullet(stripped[2:])
        elif re.match(r"^\d+\. ", stripped):
            _add_bullet(re.sub(r"^\d+\. ", "", stripped))
        elif stripped.startswith("**") and stripped.endswith("**") and stripped.count("**") == 2:
            _add_heading(stripped.strip("*"), 2)
        elif stripped.startswith(">"):
            q = re.sub(r"^(\s*>+\s?)+", "", stripped)
            p = doc.add_paragraph()
            _add_runs(p, q)
            p.paragraph_format.left_indent = Inches(0.25)
            p.paragraph_format.space_before = Pt(4)
            p.paragraph_format.space_after  = Pt(4)
        else:
            _add_body(stripped)
        i += 1

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════
#  DEEP RESEARCH (DRT) — browser-driven, agentic web research
# ══════════════════════════════════════════════════════════════

# Every tier opens with an Odysseus pre-research pass (2026-09-03) — it is the first stage.
DR_STAGES = ["odysseus", "stage1", "stage2", "stage3", "stage4", "synthesize", "report"]

# Odysseus pre-pass sizing per tier: (rounds, seconds). Runs BEFORE the research clock,
# like document analysis, so the tier's own window is untouched.
_ODY_BY_DEPTH = {"standard": (2, 150), "deep": (4, 300), "exhaustive": (6, 600)}


def _run_odysseus_prepass(job_id, job, query, clarifications, depth, stop_ev, skip_ev, prog):
    """Headless IterResearch pre-pass for every tier. Returns the report markdown ('' on
    failure/skip). STOP cancels it (the worker then honors stop_mode); the stage Skip button
    cancels it and the run continues with whatever it had. Never raises."""
    rounds, secs = _ODY_BY_DEPTH.get(depth, _ODY_BY_DEPTH["standard"])
    prog("odysseus", None, f"Odysseus pre-research — {rounds} rounds, up to {secs // 60} min…")

    def _p(ev):
        msg = "Odysseus — " + _ody_msg(ev)
        with _JOBS_LOCK:
            job["message"] = msg
        _push_event(job_id, "ody", msg)

    try:
        import asyncio
        from engines.odysseus.deep_research import DeepResearcher
        q_full = query
        if clarifications:
            q_full += f"\n\nCONTEXT / CLARIFICATIONS:\n{clarifications[:2000]}"
        from engines.research.models import get_model as _gm
        researcher = DeepResearcher(
            llm_endpoint="https://api.anthropic.com/v1/messages",  # ignored by our adapter
            llm_model=_gm("route"),      # the pre-pass is extraction/query-gen work -> Sonnet
            max_rounds=rounds, max_time=secs, progress_callback=_p)
        # Cooperative cancel: STOP or Skip-this-stage ends the pre-pass promptly.
        done = threading.Event()

        def _watch():
            while not done.wait(0.5):
                if stop_ev.is_set() or skip_ev.is_set():
                    researcher.cancel()
                    _push_event(job_id, "ody", "Odysseus pre-research cancelled by user — continuing")
                    return
        threading.Thread(target=_watch, daemon=True).start()
        try:
            report = (asyncio.run(researcher.research(q_full)) or "").strip()
        finally:
            done.set()
        if skip_ev.is_set():
            skip_ev.clear()
        if report and not report.lower().startswith("**search unavailable**"):
            _push_event(job_id, "ody",
                        f"[synth] Odysseus pre-research complete — {len(report):,} chars, "
                        f"{researcher.round_count} rounds, {len(researcher.urls_fetched)} URLs; "
                        "handing off to the browser pipeline")
            return report
        _push_event(job_id, "ody", "Odysseus pre-research produced nothing usable — continuing without it")
    except Exception as e:  # noqa: BLE001 — the pre-pass is an enhancer, never a blocker
        _push_event(job_id, "ody",
                    f"Odysseus pre-pass failed ({type(e).__name__}) — continuing with a plain run")
    return ""
_DR_EVENTS = {}     # job_id -> threading.Event (for batched Stage-4 credential prompt)
_DR_SKIP = {}       # job_id -> threading.Event (user "skip this stage" signal)
_DR_STOP = {}       # job_id -> threading.Event (user "STOP the whole run" signal)
_DR_REFUSAL = {}    # job_id -> threading.Event (user's answer to a Claude-refusal prompt)
_PROBE = {}         # probe_id -> {done, total, results:[{domain,count,ok}]} (source-access probe)
_DR_SOURCES_PATH = os.path.join(_ROOT, "config", "drt_sources.json")

_MEMO_WORD_CAP = 6000
_DOC_WORD_CAP = 50000


def _truncate_words(text: str, max_words: int) -> str:
    """Trim to the first `max_words` words, preserving formatting up to the cut."""
    if not text:
        return text
    count = 0
    for m in re.finditer(r"\S+", text):
        count += 1
        if count >= max_words:
            end = m.end()
            if end >= len(text.rstrip()):
                return text
            return text[:end] + f"\n\n[… truncated at {max_words:,} words]"
    return text


def _dr_clarify(query, provider="claude"):
    """Quick scoping pass: up to 3 clarifying questions, or none."""
    from engines.research.llm import make_client
    try:
        client = make_client(provider, os.environ.get("ANTHROPIC_API_KEY", "").strip())
    except Exception:          # LM Studio down in local mode → just skip clarification
        client = None
    if client is None:
        return {"needs_clarification": False, "questions": []}
    sys_p = ("You scope research questions for a deep web-research tool. "
             "If the question is already specific enough to research well, set "
             "needs_clarification false. Otherwise give up to 3 SHORT clarifying questions "
             "that would most sharpen the search (scope, entity, timeframe, angle). "
             "Respond with ONLY a JSON object: "
             '{"needs_clarification": boolean, "questions": ["...", "..."]}')
    from engines.research.models import get_model
    r = client.messages.create(model=get_model("route"), max_tokens=400,
                               system=sys_p, messages=[{"role": "user", "content": query}])
    txt = "".join(getattr(b, "text", "") for b in r.content if getattr(b, "type", "") == "text")
    m = re.search(r"\{.*\}", txt, re.S)
    data = json.loads(m.group(0)) if m else {}
    qs = [q for q in data.get("questions", []) if isinstance(q, str) and q.strip()][:3]
    return {"needs_clarification": bool(data.get("needs_clarification") and qs), "questions": qs}


def _dr_harvest_to_md(query, h):
    out = [f"# Research audit — sources read\n",
           f"**Question:** {query}\n",
           "> Audit trail for the synthesized report above: exactly what the 4-stage pipeline "
           "searched, read, and surfaced. Use it to trace any claim back to its page.\n",
           f"**Coverage:** {h.pages_used} pages opened from {h.searches_used} browser searches · "
           f"{len(h.items)} items harvested.\n"]
    plan = getattr(h, "plan", None)
    if isinstance(plan, dict) and plan:
        def _w(k):
            ch = plan.get(k, {});
            return f"{int(round(ch.get('weight', 0) * 100))}%" if ch.get("use", True) and ch.get("weight", 0) > 0 else "off"
        emph = plan.get("site_queries", {}).get("emphasis", [])
        et = f", favoring {', '.join(emph)}" if emph else ""
        ns = "on" if plan.get("neural_search", {}).get("use") else "off"
        out.append(
            f"**Research plan** ({'planner' if plan.get('_planned') else 'default'}): "
            f"baseline sweep {_w('api_search')} · open engines {_w('web_engines')} · "
            f"curated sites {_w('site_queries')} · API discovery (Tavily/Exa) {ns}{et}. "
            f"_{plan.get('rationale', '').strip()}_\n")
    try:
        from engines.research.models import all_models
        mm = all_models()
        out.append(
            f"**Models:** plan `{mm['plan']}` · lanes `{mm['search']}` · dig `{mm.get('dig', mm['plan'])}` · "
            f"route `{mm['route']}` · extract `{mm['extract']}` · synthesize `{mm['synthesize']}`\n")
    except Exception:
        pass
    us = getattr(h, "usage", None)
    if isinstance(us, dict) and us.get("by_model"):
        rows = " · ".join(
            f"`{m}` {u['calls']} calls, {u['input']/1000:.0f}K in + {u['cache_read']/1000:.0f}K cached"
            f" + {u['cache_write']/1000:.0f}K cache-write, {u['output']/1000:.0f}K out ≈ ${u['usd']:.2f}"
            for m, u in us["by_model"].items())
        out.append(f"**API usage this run ≈ ${us.get('total_usd', 0):.2f}** — {rows}\n")
    lanes = getattr(h, "lanes", None) or []
    if lanes:
        reps = {r.get("lane"): r.get("reason", "") for r in (getattr(h, "lane_reports", None) or [])}
        out.append(f"**Parallel research lanes ({len(lanes)}):**\n")
        for i, ln in enumerate(lanes, 1):
            fin = reps.get(f"L{i}", "")
            out.append(f"- **L{i} · {ln.get('name', '')}** — {ln.get('mission', '')}"
                       + (f"\n  _finished: {fin[:300]}_" if fin else ""))
        if reps.get("dig"):
            out.append(f"- **Hot-trail dig** — _finished: {reps['dig'][:300]}_")
        out.append("")
    fbs = getattr(h, "fallback_fetches", None) or []
    if fbs:
        out.append(f"**Pages Chrome could not read, fetched via unblocking providers ({len(fbs)}):** "
                   + "; ".join(f"{f.get('url', '')[:70]} ({f.get('via')})" for f in fbs[:12]) + "\n")
    if getattr(h, "curated_searched", None):
        out.append(f"**Curated sites searched ({len(h.curated_searched)}):** "
                   f"{', '.join(h.curated_searched)}\n")
    if getattr(h, "exa_searches", 0) or getattr(h, "exa_similar", 0):
        out.append(f"**Neural search (Exa):** {h.exa_searches} queries · "
                   f"{h.exa_similar} find-similar · "
                   f"{len(getattr(h, 'exa_urls', []))} URLs surfaced\n")
    if getattr(h, "logged_in", None):
        out.append(f"**Logged in this run:** {', '.join(h.logged_in)}\n")
    if getattr(h, "login_warnings", None):
        out.append("**⚠ Stored-login warnings:** "
                   + "; ".join(f"{w.get('domain')} — {w.get('detail')}" for w in h.login_warnings)
                   + "\n")
    if getattr(h, "skipped_gated", None):
        out.append(f"**Gated sources skipped** (login failed / needed 2FA): "
                   f"{', '.join(h.skipped_gated)}\n")
    if getattr(h, "gated_candidates", None):
        out.append(f"**Gated sources not pursued:** {', '.join(h.gated_candidates)}\n")
    if h.stopped_reason:
        out.append(f"**Stop reason:** {h.stopped_reason}\n")

    stage1 = next((it for it in h.items if it.via == "stage1"), None)
    if stage1 and stage1.text.strip():
        out.append("\n## Stage 1 — baseline browser brief\n")
        out.append(stage1.text.strip())

    if h.agent_notes:
        out.append("\n## Agent notes (Stages 2–4)\n")
        out.append(h.agent_notes)

    opened = [it for it in h.items if it.via != "stage1"]
    out.append("\n## Pages read (Stages 2–4)\n")
    if not opened:
        out.append("_No pages opened — nothing of value surfaced beyond the Stage 1 brief._")
    else:
        exa_set = set(getattr(h, "exa_urls", []) or [])
        for i, it in enumerate(opened, 1):
            title = (it.title or it.url).strip()
            flag = " · ⚠ thin/screenshot" if it.used_screenshot else ""
            exa_flag = " · 🔮 via Exa" if it.url in exa_set else ""
            out.append(f"{i}. [{title}]({it.url}) — *{it.via}/{it.source_type}*, "
                       f"{len(it.text)} chars{flag}{exa_flag}")

    if getattr(h, "stage1_sources", None):
        out.append("\n## Other sources surfaced by the baseline sweep (not opened)\n")
        for s in h.stage1_sources[:25]:
            out.append(f"- [{(s.get('title') or s.get('url'))}]({s.get('url')})")
    return "\n".join(out)


def _dr_wrap_report(query, h, synth_md, ody_md=""):
    """Wrap the synthesized report with the harvest audit trail (collapsed <details> for
    the web UI; a plain section for the docx). `ody_md` = the Odysseus pre-research pass
    that grounded the run — appended as its own collapsed/trailing section for traceability."""
    harvest_md = _dr_harvest_to_md(query, h)
    if (ody_md or "").strip():
        harvest_md += ("\n\n## Odysseus pre-research (headless pass that grounded this run)\n\n"
                       + ody_md.strip())
    synth_md = (synth_md or "").strip()
    warn = ""
    if getattr(h, "login_warnings", None):
        lines = "; ".join(f"**{w.get('domain')}** ({w.get('detail')})" for w in h.login_warnings)
        warn = ("> ⚠ **Stored-login warning:** the tool could not complete the saved login for "
                + lines + ". Results from those sites may be limited to publicly visible content — "
                "re-check the credentials via the 🔑 lock in the Sources panel.\n\n")
    if not synth_md:
        return warn + harvest_md, warn + harvest_md
    display_md = (warn + synth_md
                  + "\n\n<details>\n<summary><strong>Research audit — what the tool read"
                    " (harvest detail)</strong></summary>\n\n"
                  + harvest_md + "\n\n</details>")
    docx_md = warn + synth_md + "\n\n---\n\n" + harvest_md
    return display_md, docx_md


_DOC_ANALYSIS_CHARS = 40000   # per-file slice fed to the analyst (full text still reaches synthesis)


def _analyze_docs(client, query, doc_parts, log=lambda m: None):
    """Distill each uploaded document against the research question — one call per file
    through the wrapped client (Fable 5 at the standard effort). Returns a compact
    combined brief used to ground planning + searching; a file whose analysis fails
    degrades to a raw excerpt, and the FULL raw text still reaches synthesis either way."""
    from engines.research.models import get_model
    briefs = []
    sys_p = (
        "You are the document-intake analyst of a deep web-research tool. The user attached "
        "a file as context for their research question. Distill THIS file against the "
        "question into a terse brief:\n"
        "- KEY FACTS / CLAIMS in the file that bear on the question (tight bullets; keep "
        "specifics — names, numbers, dates; attribute claims as the file does).\n"
        "- ENTITIES & TERMS worth searching (including any insider vocabulary the file uses).\n"
        "- OPEN QUESTIONS the file raises or leaves unanswered.\n"
        "- SEARCH LEADS (2-4) the file suggests pursuing on the live web.\n"
        "Signal only, no filler. If the file has nothing relevant to the question, say so "
        "in one line.")
    for part in doc_parts or []:
        head, _, body = part.partition("\n")
        name = head.strip().strip("[]") or "document"
        if not body.strip():
            continue
        log(f"[docs] analyzing {name} ({len(body):,} chars)…")
        try:
            r = client.messages.create(
                model=get_model("route"), max_tokens=1500, system=sys_p,
                messages=[{"role": "user", "content":
                           f"RESEARCH QUESTION:\n{query}\n\nFILE: {name}\n\n"
                           f"FILE CONTENTS:\n{body[:_DOC_ANALYSIS_CHARS]}"}])
            brief = "".join(getattr(b, "text", "") for b in r.content
                            if getattr(b, "type", "") == "text").strip()
            if brief:
                briefs.append(f"### {name}\n{brief}")
                log(f"[docs] ✓ {name} distilled → {len(brief):,} chars of brief")
                continue
        except Exception as e:  # noqa: BLE001
            log(f"[docs] ✗ analysis failed for {name} ({type(e).__name__}) — using raw excerpt")
        briefs.append(f"### {name} (raw excerpt — analysis unavailable)\n{body[:2000]}")
    return "\n\n".join(briefs)


def _dr_worker(job_id, query, depth, clarifications, doc_context, channel_overrides=None,
               provider="claude", doc_parts=None):
    job = _JOBS[job_id]
    ev = _DR_EVENTS[job_id]
    skip_ev = _DR_SKIP[job_id]
    stop_ev = _DR_STOP[job_id]

    def prog(stage, pct, message):
        with _JOBS_LOCK:
            job["stage"] = stage
            job["pct"] = pct
            job["message"] = message
        _push_event(job_id, "stage", message)

    log_fn = lambda m: _push_event(job_id, _event_kind(m), m)  # noqa: E731

    def request_credentials(candidates):
        ev.clear()
        with _JOBS_LOCK:
            job["awaiting_credentials"] = candidates
            job["submitted_credentials"] = None
            job["message"] = f"Waiting for credentials for {len(candidates)} gated source(s)…"
        ok = ev.wait(timeout=600)
        with _JOBS_LOCK:
            submitted = job.get("submitted_credentials") or {}
            job["awaiting_credentials"] = None
            job["submitted_credentials"] = None
        return submitted if ok else {}

    try:
        from engines.research.agent import run_search
        from engines.research.browser import DRTBrowser
        from engines.research.agent import _load_governance, run_gap_round
        from engines.research.synthesize import (synthesize, classify_category,
                                                 stop_judge, gap_queries, DEEPEN_ROUNDS)
        from engines.research.llm import make_client, is_local, LocalLLMUnavailable
        gov = _load_governance()
        api_key = os.environ.get("ANTHROPIC_API_KEY", "").strip()
        # Refusal hook: if Claude declines (safety refusal), pause the pipeline right
        # there, flag the UI, and wait (up to 10 min) for the user's choice. "local"
        # → the wrapper swaps to LM Studio and RESUMES from the refused call; "abort"
        # → stop gathering and assemble what we have. Reuses the local-model wiring.
        refusal_ev = _DR_REFUSAL[job_id]

        def _refusal_hook(details):
            # Pause: flag the UI, block for the user's decision. UI-only — the wrapper
            # decides whether to stop (via _on_stop) based on the actual swap result.
            refusal_ev.clear()
            with _JOBS_LOCK:
                job["awaiting_refusal_choice"] = details
                job["refusal_decision"] = None
            log_fn(f"[refusal] Claude declined ({details.get('category')}) — "
                   f"awaiting your choice (switch to local model / stop)")
            ok = refusal_ev.wait(timeout=600)
            with _JOBS_LOCK:
                dec = (job.get("refusal_decision") or "abort") if ok else "abort"
                job["awaiting_refusal_choice"] = None
            log_fn("[refusal] switching to the local model and resuming…" if dec == "local"
                   else "[refusal] stopping — assembling from what was gathered")
            return dec

        def _refusal_stop():
            # The wrapper calls this only when the run must halt (abort, or a local
            # swap that failed). Assemble a report from what was already gathered.
            stop_ev.set()
            with _JOBS_LOCK:
                job["stop_mode"] = "assemble"

        if is_local(provider):
            try:
                client = make_client("local")           # LMStudioClient (whatever model is loaded)
            except LocalLLMUnavailable as e:
                with _JOBS_LOCK:
                    job["error"] = str(e); job["done"] = True
                return
        else:
            client = make_client("claude", api_key, refusal_hook=_refusal_hook,
                                 on_stop=_refusal_stop)

        # ── Odysseus pre-research — EVERY tier (2026-09-03). The headless IterResearch
        # pass runs first, sized to the tier, and its report is fed in EXACTLY like a
        # dropped file: the intake analyst distills it into grounding for the plan +
        # lanes, and its full text reaches synthesis as trusted user material. Local-model
        # runs skip it (Odysseus is wired to the Anthropic key).
        from engines.research.agent import normalize_depth as _nd0
        depth = _nd0(depth)
        ody_report = ""
        if not is_local(provider) and not stop_ev.is_set():
            ody_report = _run_odysseus_prepass(job_id, job, query, clarifications, depth,
                                               stop_ev, skip_ev, prog)
            if stop_ev.is_set() and (job.get("stop_mode") == "abort"):
                with _JOBS_LOCK:
                    job["aborted"] = True; job["message"] = "Stopped — discarded."
                    job["done"] = True
                return
        else:
            log_fn("[ody] Odysseus pre-research skipped (local model run)" if is_local(provider)
                   else "[ody] Odysseus pre-research skipped")
        if ody_report:
            ody_part = ("[Odysseus pre-research findings]\n"
                        + _truncate_words(ody_report, _DOC_WORD_CAP))
            doc_parts = list(doc_parts or []) + [ody_part]
            doc_context = (doc_context + "\n\n" + ody_part).strip() if doc_context else ody_part

        # ── Uploaded documents: extract → parse (done at POST) → ANALYZE here.
        # Each file is distilled against the question (facts, entities, open questions,
        # search leads); the brief grounds the planner + browser agent. The FULL raw
        # text still goes to synthesis as trusted user_docs. Runs before the research
        # clock starts, so analysis never eats the gathering window.
        clar = clarifications or ""
        if doc_context:
            brief = ""
            if client and doc_parts:
                prog("stage1", None, f"Analyzing {len(doc_parts)} uploaded document(s)…")
                brief = _analyze_docs(client, query, doc_parts, log=log_fn)
            if brief:
                clar = (clar + "\n\nSUPPORTING DOCUMENT ANALYSIS (distilled from the user's "
                        "uploaded files; their full text is applied again at synthesis):\n"
                        + brief).strip()
            else:
                clar = (clar + "\n\nSUPPORTING DOCUMENT EXCERPTS (user-provided):\n"
                        + doc_context).strip()

        prog("stage1", None, "Starting…")
        br = DRTBrowser(log=log_fn).start()
        report_md_synth = ""
        synth_error = None
        try:
            # TIME-boxed run: the tier defines the whole run's window; gathering gets
            # the window minus the synthesis reserve, deepening only runs on leftover time.
            import time as _time
            from engines.research.agent import DEPTH_BUDGETS as _DB, normalize_depth as _nd
            depth = _nd(depth)
            _b = _DB[depth]
            run_deadline = _time.time() + _b["seconds"]
            h = run_search(query, depth=depth, clarifications=clar, browser=br,
                           progress=prog, log=log_fn,
                           request_credentials=request_credentials, skip_event=skip_ev,
                           channel_overrides=channel_overrides, client=client, provider=provider,
                           deadline=run_deadline - _b["reserve"], stop_event=stop_ev)
            # User pressed STOP and chose "start a new query" → discard, no report.
            if stop_ev.is_set() and (job.get("stop_mode") == "abort"):
                with _JOBS_LOCK:
                    job["aborted"] = True; job["message"] = "Stopped — discarded."
                    job["done"] = True
                return
            if stop_ev.is_set():
                log_fn("[stop] assembling report from what was gathered before stop…")
            if client and (h.items or doc_context):
                try:
                    prog("synthesize", None, "Synthesizing the report…")
                    category = getattr(h, "category", "") or classify_category(client, query)
                    cache = {}
                    synth = synthesize(h, gov, client, progress=prog, nugget_cache=cache,
                                       category=category, user_docs=doc_context, log=log_fn)
                    report_md_synth = synth.get("report_md", "")
                    h.category = category
                    extra = DEEPEN_ROUNDS.get(depth, 1)
                    for rnd in range(1, extra + 1):
                        if stop_ev.is_set():
                            break     # user stopped — no extra gathering, just synthesize
                        if skip_ev.is_set():
                            skip_ev.clear(); break
                        if _time.time() > run_deadline - 75:
                            log_fn("[synth] time window spent — skipping further deepening")
                            break
                        stop, reason = stop_judge(client, query, report_md_synth, log=log_fn)
                        if stop:
                            h.stopped_reason = reason or "report judged comprehensive"
                            break
                        gaps = gap_queries(client, query, report_md_synth, log=log_fn)
                        if not gaps:
                            break
                        prog("synthesize", None, f"Deepening (round {rnd}): {gaps[0][:60]}")
                        added = run_gap_round(client, br, h, gaps, governance=gov,
                                              progress=prog, log=log_fn, skip_event=skip_ev,
                                              deadline=run_deadline - 75, stop_event=stop_ev)
                        if not added:
                            break
                        synth = synthesize(h, gov, client, progress=prog, nugget_cache=cache,
                                           category=category, user_docs=doc_context, log=log_fn)
                        report_md_synth = synth.get("report_md", "")
                except Exception as e:  # noqa: BLE001 — preserve the harvest
                    synth_error = e
                    prog("synthesize", None, "Synthesis interrupted — assembling harvested findings…")
        finally:
            try:
                br.close()
            except Exception:
                pass

        if synth_error is not None:
            detail = str(synth_error).strip()
            if "credit balance is too low" in detail.lower():
                hint = ("the Anthropic API account is **out of credits** — top up at "
                        "console.anthropic.com → Plans & Billing, then re-run")
            else:
                hint = "re-run once the API is reachable"
            report_md_synth = (
                f"> ⚠️ **Report synthesis was interrupted** ({type(synth_error).__name__}); "
                f"{hint}. The web harvest completed and the gathered sources are preserved below.\n>\n"
                f"> _Detail: {detail[:300]}_\n\n"
                + (report_md_synth or "")
            )

        # "How to go deeper" assessment — appended to the report (display + docx)
        # and surfaced separately for the UI's framed panel. Best-effort.
        go_deeper = ""
        if client and report_md_synth and synth_error is None:
            try:
                from engines.research.synthesize import deeper_assessment
                go_deeper = (deeper_assessment(client, query, report_md_synth, h,
                                               log=log_fn) or "").strip()
            except Exception:
                go_deeper = ""
            if go_deeper:
                report_md_synth = report_md_synth + "\n\n" + go_deeper

        prog("report", None, "Assembling report…")
        try:      # real spend for the audit (ClaudeClient counts every call's usage)
            h.usage = client.usage_summary() if hasattr(client, "usage_summary") else None
            if h.usage:
                log_fn(f"[synth] API usage this run ≈ ${h.usage['total_usd']:.2f} "
                       f"({h.usage['totals']['input']//1000}K in, "
                       f"{h.usage['totals']['cache_read']//1000}K cached, "
                       f"{h.usage['totals']['output']//1000}K out)")
        except Exception:
            h.usage = None
        report_md, docx_md = _dr_wrap_report(query, h, report_md_synth, ody_md=ody_report)
        try:
            docx_b64 = base64.b64encode(_memo_to_docx_bytes(docx_md, "Deep Research")).decode()
        except Exception:
            docx_b64 = ""
        # Subject-appropriate document title (one cheap call) → the saved filename,
        # the Save & open button, and the results header.
        title = _report_title(query, report_md_synth, client) if client else _report_title(query)
        log_fn(f"[synth] title: {title}")
        sources = [{"title": it.title, "url": it.url, "type": it.source_type,
                    "chars": len(it.text)} for it in h.items if it.via != "stage1"]
        result = {"query": query, "title": title, "report_md": report_md, "sources": sources,
                  "source_count": len(sources), "docx_b64": docx_b64,
                  "saved_path": _autosave_report(docx_b64, query, "Deep Research", title=title),
                  "go_deeper_md": go_deeper,
                  "odysseus_report_md": ody_report,
                  "discovered_forums": getattr(h, "discovered_forums", []) or [],
                  "category": getattr(h, "category", ""),
                  "plan": getattr(h, "plan", {}),
                  "stats": {"searches": h.searches_used, "pages": h.pages_used,
                            "category": getattr(h, "category", "") or "general",
                            "stopped": h.stopped_reason}}
        with _JOBS_LOCK:
            job["result"] = result
            job["stage"] = "report"; job["pct"] = 100; job["message"] = "Done"
            job["done"] = True
    except Exception:
        with _JOBS_LOCK:
            job["error"] = traceback.format_exc()
            job["done"] = True


@app.route("/api/deep_research", methods=["POST", "OPTIONS"])
def deep_research_start():
    if request.method == "OPTIONS":
        return "", 204
    query = (request.form.get("query") or "").strip()
    if not query:
        return jsonify({"error": "query is required"}), 400
    from engines.research.agent import normalize_depth
    raw_depth = (request.form.get("depth") or "standard").strip().lower()
    depth = normalize_depth(raw_depth)
    clarifications = (request.form.get("clarifications") or "").strip()
    provider = (request.form.get("provider") or "claude").strip().lower()

    refinement = (request.form.get("refinement") or "").strip()
    prior_report = (request.form.get("prior_report") or "").strip()
    if refinement:
        extra = ("\n\nThis is a REFINEMENT of a prior research run. Focus on the refinement; "
                 "build on what is already known and do NOT repeat it.\n"
                 f"REFINEMENT REQUEST:\n{refinement}")
        if prior_report:
            extra += f"\n\nPRIOR REPORT (context — already produced):\n{prior_report[:6000]}"
        clarifications = (clarifications + extra).strip()

    memo_filename = (request.form.get("memo_filename") or "").strip()

    channel_overrides = {}
    try:
        raw = json.loads(request.form.get("channels") or "{}")
        if isinstance(raw, dict):
            channel_overrides = {k: bool(raw[k]) for k in
                                 ("api_search", "web_engines", "site_queries", "neural_search") if k in raw}
    except Exception:
        channel_overrides = {}

    doc_parts = []
    for f in request.files.getlist("files"):
        if not f or not f.filename:
            continue
        ext = f.filename.rsplit(".", 1)[-1].lower() if "." in f.filename else "tmp"
        tmp = tempfile.NamedTemporaryFile(suffix="." + ext, delete=False)
        try:
            f.save(tmp.name); tmp.close()
            txt = _extract_file_text(tmp.name, f.filename)
            if txt:
                cap = _MEMO_WORD_CAP if (memo_filename and f.filename == memo_filename) else _DOC_WORD_CAP
                doc_parts.append(f"[{f.filename}]\n{_truncate_words(txt, cap)}")
        except Exception as e:  # noqa: BLE001
            doc_parts.append(f"[{f.filename}] (could not read: {e})")
        finally:
            try:
                os.unlink(tmp.name)
            except Exception:
                pass
    doc_context = "\n\n".join(doc_parts)

    # (The old "Deep + Odysseus" chained depth is gone as a separate path — every tier
    # now opens with an Odysseus pre-pass inside _dr_worker; normalize_depth maps the
    # legacy "odysseus" value to the deep tier.)
    job_id = os.urandom(8).hex()
    with _JOBS_LOCK:
        _JOBS[job_id] = {
            "stage": "odysseus", "pct": None, "message": "Starting…",
            "done": False, "error": None, "result": None,
            "awaiting_credentials": None, "submitted_credentials": None,
            "events": [], "eseq": 0,
        }
    _DR_EVENTS[job_id] = threading.Event()
    _DR_SKIP[job_id] = threading.Event()
    _DR_STOP[job_id] = threading.Event()
    _DR_REFUSAL[job_id] = threading.Event()
    threading.Thread(target=_dr_worker,
                     args=(job_id, query, depth, clarifications, doc_context, channel_overrides,
                           provider, doc_parts),
                     daemon=True).start()
    return jsonify({"job_id": job_id, "stages": DR_STAGES}), 202


@app.route("/api/deep_research/status", methods=["GET"])
def deep_research_status():
    job_id = request.args.get("job", "")
    try:
        after = int(request.args.get("after", 0) or 0)
    except (TypeError, ValueError):
        after = 0
    with _JOBS_LOCK:
        job = _JOBS.get(job_id)
        if not job:
            return jsonify({"error": "unknown job"}), 404
        payload = {
            "stage": job["stage"], "pct": job["pct"], "message": job["message"],
            "done": job["done"], "error": job["error"], "stages": DR_STAGES,
            "awaiting_credentials": job.get("awaiting_credentials"),
            "awaiting_refusal_choice": job.get("awaiting_refusal_choice"),
            "events": [e for e in (job.get("events") or []) if e[0] > after],
            "eseq": job.get("eseq", 0),
        }
        payload["aborted"] = bool(job.get("aborted"))
        if job["done"] and not job["error"]:
            payload["result"] = job["result"]
        if job["done"]:
            _JOBS.pop(job_id, None)
            _DR_EVENTS.pop(job_id, None)
            _DR_SKIP.pop(job_id, None)
            _DR_STOP.pop(job_id, None)
            _DR_REFUSAL.pop(job_id, None)
    return jsonify(payload)


@app.route("/api/deep_research/skip_stage", methods=["POST", "OPTIONS"])
def deep_research_skip_stage():
    """End the current browser stage early (Stages 2–4). The pipeline rolls forward."""
    if request.method == "OPTIONS":
        return "", 204
    job_id = request.args.get("job", "") or (request.form.get("job") or "")
    ev = _DR_SKIP.get(job_id)
    if ev:
        ev.set()
    return jsonify({"ok": True})


@app.route("/api/deep_research/stop", methods=["POST", "OPTIONS"])
def deep_research_stop():
    """STOP the whole run. Body/query `mode`:
      • 'assemble' (default) — halt gathering, synthesize a report from what's harvested.
      • 'abort' — halt gathering and discard (no report; the UI returns to a new query).
    Also unblocks a Stage-4 credential wait so the worker doesn't sit for 10 min."""
    if request.method == "OPTIONS":
        return "", 204
    job_id = request.args.get("job", "") or (request.form.get("job") or "")
    data = request.get_json(silent=True) or {}
    mode = (data.get("mode") or request.args.get("mode") or "assemble").strip().lower()
    if mode not in ("assemble", "abort"):
        mode = "assemble"
    with _JOBS_LOCK:
        job = _JOBS.get(job_id)
        if job is not None:
            job["stop_mode"] = mode
    ev = _DR_STOP.get(job_id)
    if ev:
        ev.set()
    skip = _DR_SKIP.get(job_id)     # break the current stage promptly too
    if skip:
        skip.set()
    cred = _DR_EVENTS.get(job_id)   # unblock a pending Stage-4 credential wait
    if cred:
        cred.set()
    return jsonify({"ok": True, "mode": mode})


@app.route("/api/deep_research/refusal_choice", methods=["POST", "OPTIONS"])
def deep_research_refusal_choice():
    """Answer a Claude-refusal prompt. Body/query `mode`:
      • 'local' — switch to the local model and RESUME from the refused call (the UI
        confirms a local model is loaded, via /local_model, before sending this).
      • 'abort' — stop and assemble a report from what was already gathered.
    Unblocks the paused worker thread."""
    if request.method == "OPTIONS":
        return "", 204
    job_id = request.args.get("job", "") or (request.form.get("job") or "")
    data = request.get_json(silent=True) or {}
    mode = (data.get("mode") or request.args.get("mode") or "abort").strip().lower()
    if mode not in ("local", "abort"):
        mode = "abort"
    with _JOBS_LOCK:
        job = _JOBS.get(job_id)
        if job is not None:
            job["refusal_decision"] = mode
    ev = _DR_REFUSAL.get(job_id)
    if ev:
        ev.set()
    return jsonify({"ok": True, "mode": mode})


# ── source-access self-probe ────────────────────────────────────────────────
# Runs each site's PUBLIC (logged-out) native search from THIS server's own browser
# environment. On the hosted instance that is a headless bundled Chromium on Render's
# datacenter IP — i.e. ground truth for "will this site's search work on centralindustrial.ai
# without a login", which browser stealth alone can't guarantee against datacenter-IP blocks.
_PROBE_QUERIES = {
    # a benign, likely-to-hit query per domain; default used for anything not listed
    "avforums.com": "receiver", "rivianforums.com": "update", "wilderssecurity.com": "firewall",
    "forum.openwrt.org": "firmware", "forum.wiimhome.com": "firmware", "github.com": "cli",
    "reddit.com": "review", "seekingalpha.com": "earnings", "x.com": "news",
}


def _probe_worker(probe_id, domains):
    import shutil
    import tempfile
    from engines.research.browser import DRTBrowser
    tmp = tempfile.mkdtemp(prefix="drt-probe-")
    try:
        # Guest context (throwaway profile, no vault) + headless — matches what the hosted
        # instance can do for a logged-out visitor.
        br = DRTBrowser(profile_dir=tmp, headed=False, log=lambda m: None).start()
        try:
            for dom in domains:
                q = _PROBE_QUERIES.get(dom, "guide")
                try:
                    n = len(br.native_search(dom, q, limit=5))
                except Exception:  # noqa: BLE001
                    n = -1
                with _JOBS_LOCK:
                    p = _PROBE.get(probe_id)
                    if p is None:
                        break
                    p["results"].append({"domain": dom, "count": n, "ok": n > 0})
        finally:
            try:
                br.close()
            except Exception:
                pass
    except Exception as e:  # noqa: BLE001
        with _JOBS_LOCK:
            if _PROBE.get(probe_id) is not None:
                _PROBE[probe_id]["error"] = f"{type(e).__name__}: {e}"
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
        with _JOBS_LOCK:
            if _PROBE.get(probe_id) is not None:
                _PROBE[probe_id]["done"] = True


@app.route("/api/deep_research/probe_sources", methods=["POST", "OPTIONS"])
def deep_research_probe_sources():
    """Start a source-access probe from THIS server. Body: {domains?: [...]}. Defaults to the
    catalog's domains. Returns {probe_id}; poll /probe_status."""
    if request.method == "OPTIONS":
        return "", 204
    data = request.get_json(silent=True) or {}
    domains = data.get("domains")
    if not isinstance(domains, list) or not domains:
        from engines.research.agent import load_sources
        domains = [(s.get("domain") or "").strip().lower()
                   for s in load_sources() if (s.get("domain") or "").strip()]
    # de-dup, cap, sanitize
    seen, clean = set(), []
    for d in domains:
        d = str(d).strip().lower()
        if d and d not in seen and re.match(r"^[a-z0-9.\-]+\.[a-z]{2,}$", d):
            seen.add(d)
            clean.append(d)
        if len(clean) >= 40:
            break
    probe_id = os.urandom(6).hex()
    with _JOBS_LOCK:
        _PROBE[probe_id] = {"done": False, "total": len(clean), "results": [], "error": None}
    threading.Thread(target=_probe_worker, args=(probe_id, clean), daemon=True).start()
    return jsonify({"probe_id": probe_id, "total": len(clean)}), 202


# ── publish the site list to the hosted instance (LOCAL-ONLY) ───────────────────────
# The hosted site gets its site list from the COMMITTED config file baked into the container
# at deploy. So "publish" = copy the local drt_sources.json into a persistent clone of the
# repo, commit, push → Render auto-deploys (~2 min). Only the site list travels; the vault
# never does. Refused on the hosted instance (it is the destination, not the source).
_REPO_DIR = os.environ.get("DRT_REPO_DIR", r"D:\_______Claude\Central-Industrial")
_REPO_URL = "https://github.com/nurtrino/Central-Industrial.git"
_REPO_SOURCES_REL = "deep-research/config/drt_sources.json"


def _is_hosted_instance() -> bool:
    # The auth gate is only configured on the hosted deployment.
    return bool(os.environ.get("AUTH_SECRET", "").strip())


def _sources_hash() -> str:
    """Short content hash of the site list — exposed on /api/health so a publish can be
    confirmed to have landed on the hosted instance (compare local vs hosted)."""
    import hashlib
    try:
        with open(_DR_SOURCES_PATH, "rb") as fh:
            return hashlib.sha1(fh.read()).hexdigest()[:8]
    except Exception:
        return ""


@app.route("/api/deep_research/publish_sources", methods=["POST", "OPTIONS"])
def deep_research_publish_sources():
    if request.method == "OPTIONS":
        return "", 204
    if _is_hosted_instance():
        return jsonify({"error": "Publishing is local-only — the hosted site is the destination, "
                                 "not the source. Edit the list on the local app and publish from there."}), 403
    import shutil
    import subprocess

    def git(*args, cwd=None):
        r = subprocess.run(["git", *args], cwd=cwd or _REPO_DIR, capture_output=True,
                           text=True, timeout=240)
        return r.returncode, ((r.stdout or "") + (r.stderr or "")).strip()

    try:
        if not os.path.isdir(os.path.join(_REPO_DIR, ".git")):
            os.makedirs(os.path.dirname(_REPO_DIR), exist_ok=True)
            rc, out = git("clone", _REPO_URL, _REPO_DIR, cwd=os.path.dirname(_REPO_DIR))
            if rc != 0:
                return jsonify({"error": f"clone failed: {out[-400:]}"}), 500
        rc, out = git("pull", "--ff-only", "origin", "main")
        if rc != 0:
            return jsonify({"error": f"pull failed: {out[-400:]}"}), 500
        dest = os.path.join(_REPO_DIR, *_REPO_SOURCES_REL.split("/"))
        os.makedirs(os.path.dirname(dest), exist_ok=True)
        shutil.copyfile(_DR_SOURCES_PATH, dest)
        git("add", _REPO_SOURCES_REL)                       # stage ONLY the site list
        rc, _ = git("diff", "--cached", "--quiet", "--", _REPO_SOURCES_REL)
        if rc == 0:
            return jsonify({"ok": True, "changed": False, "sources_hash": _sources_hash(),
                            "message": "Hosted site list already matches — nothing to publish."})
        from engines.research.agent import load_sources
        n = len(load_sources())
        rc, out = git("commit", "-m", f"Deep Research: site list published from local ({n} sites)")
        if rc != 0:
            return jsonify({"error": f"commit failed: {out[-400:]}"}), 500
        rc, out = git("push", "origin", "main")
        if rc != 0:
            return jsonify({"error": f"push failed: {out[-400:]}"}), 500
        _, sha = git("rev-parse", "--short", "HEAD")
        return jsonify({"ok": True, "changed": True, "commit": sha, "sources_hash": _sources_hash(),
                        "message": f"Published {n} sites (commit {sha}). Render is redeploying the "
                                   f"hosted site now — about 2 minutes."})
    except Exception as e:  # noqa: BLE001
        return jsonify({"error": f"{type(e).__name__}: {e}"}), 500


@app.route("/api/deep_research/probe_status", methods=["GET"])
def deep_research_probe_status():
    pid = request.args.get("probe", "")
    with _JOBS_LOCK:
        p = _PROBE.get(pid)
        if not p:
            return jsonify({"error": "unknown probe"}), 404
        payload = {"done": p["done"], "total": p["total"],
                   "results": list(p["results"]), "error": p.get("error")}
        if p["done"]:
            _PROBE.pop(pid, None)
    return jsonify(payload)


@app.route("/api/deep_research/credentials", methods=["POST", "OPTIONS"])
def deep_research_credentials():
    """Stage-4 credential submission (batched). Body: {credentials: {domain: {username,
    password, login_url} | null}}. Unblocks the worker; null/missing entries are skipped."""
    if request.method == "OPTIONS":
        return "", 204
    job_id = request.args.get("job", "") or (request.form.get("job") or "")
    data = request.get_json(silent=True) or {}
    creds = data.get("credentials", {}) if isinstance(data, dict) else {}
    with _JOBS_LOCK:
        job = _JOBS.get(job_id)
        if job is not None:
            job["submitted_credentials"] = creds
    ev = _DR_EVENTS.get(job_id)
    if ev:
        ev.set()
    return jsonify({"ok": True})


@app.route("/api/deep_research/clarify", methods=["POST", "OPTIONS"])
def deep_research_clarify():
    if request.method == "OPTIONS":
        return "", 204
    query = (request.form.get("query") or "").strip()
    provider = (request.form.get("provider") or "claude").strip().lower()
    if not query:
        return jsonify({"needs_clarification": False, "questions": []})
    try:
        return jsonify(_dr_clarify(query, provider))
    except Exception as e:  # noqa: BLE001
        return jsonify({"needs_clarification": False, "questions": [], "warn": str(e)})


@app.route("/api/deep_research/local_model", methods=["GET"])
def deep_research_local_model():
    """Report the model currently loaded in LM Studio (for the Claude/Local-AI switch)."""
    from engines.research.llm import detect_local_model, LMSTUDIO_URL
    try:
        return jsonify({"ok": True, "model": detect_local_model(), "url": LMSTUDIO_URL})
    except Exception as e:  # noqa: BLE001
        return jsonify({"ok": False, "error": str(e), "url": LMSTUDIO_URL})


@app.route("/api/deep_research/vault", methods=["GET", "POST", "OPTIONS"])
def deep_research_vault():
    """Manage stored credentials (encrypted local vault).
    GET -> {"domains": [...]}   POST {domain, username, password, login_url} -> store
    POST {domain, delete: true} -> remove."""
    if request.method == "OPTIONS":
        return "", 204
    from engines.research.login import CredentialVault, normalize_domain
    v = CredentialVault()
    if request.method == "GET":
        return jsonify({"domains": v.domains()})
    data = request.get_json(silent=True) or {}
    domain = normalize_domain(data.get("domain") or "")
    if not domain:
        return jsonify({"error": "domain required"}), 400
    if data.get("delete"):
        sites = v.load(); sites.pop(domain, None); v.save(sites)
        return jsonify({"ok": True, "deleted": domain, "domains": v.domains()})
    if not (data.get("username") or "").strip():
        return jsonify({"error": "username required"}), 400
    v.set(domain, data.get("username", ""), data.get("password", ""),
          login_url=data.get("login_url", ""))
    return jsonify({"ok": True, "domains": v.domains()})


@app.route("/api/deep_research/sources", methods=["GET", "POST", "OPTIONS"])
def deep_research_sources():
    if request.method == "OPTIONS":
        return "", 204
    if request.method == "GET":
        try:
            with open(_DR_SOURCES_PATH, encoding="utf-8") as fh:
                return jsonify(json.load(fh))
        except Exception:
            return jsonify({"sources": []})
    data = request.get_json(silent=True) or {}
    srcs = data.get("sources", [])
    from engines.research.login import normalize_domain
    # Sanitize to the UNIFIED schema: {url, login_required, note?, search_url?}. Accepts legacy
    # entries (domain/description) and drops the retired per-site fields (name/type/enabled).
    clean, seen = [], set()
    for s in srcs:
        if not isinstance(s, dict):
            continue
        url = (s.get("url") or s.get("domain") or "").strip()
        dom = normalize_domain(url) or url.lower()
        if not dom or dom in seen:
            continue
        seen.add(dom)
        lr = s.get("login_required")
        entry = {"url": dom, "login_required": lr if lr in (True, False) else None}
        note = (s.get("note") or s.get("description") or "").strip()
        if note:
            entry["note"] = note
        surl = (s.get("search_url") or "").strip()
        if surl:
            entry["search_url"] = surl
        clean.append(entry)
    os.makedirs(os.path.dirname(_DR_SOURCES_PATH), exist_ok=True)
    with open(_DR_SOURCES_PATH, "w", encoding="utf-8") as fh:
        json.dump({"_comment": ("Unified site list: url + login_required (green/red dot). "
                                "false = searchable logged out; true = needs a login (vault locally, "
                                "one-time prompt otherwise); null = untested. The research plan selects "
                                "the topically-relevant sites per query."),
                   "sources": clean}, fh, indent=2)
    return jsonify({"ok": True, "count": len(clean)})


# ── Extract a source: scrape one page / crawl a site (Firecrawl → Exa fallback) ──
def _firecrawl_scrape(url, log=lambda m: None):
    """Single-page scrape via Firecrawl v2. Returns dict or None."""
    key = os.environ.get("FIRECRAWL_API_KEY", "").strip()
    if not key:
        return None
    import requests
    try:
        r = requests.post("https://api.firecrawl.dev/v2/scrape",
                          headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
                          json={"url": url, "formats": ["markdown"]}, timeout=120)
        if r.status_code == 200:
            data = (r.json() or {}).get("data") or {}
            md = (data.get("markdown") or "").strip()
            if md:
                meta = data.get("metadata") or {}
                return {"markdown": md, "title": meta.get("title") or url,
                        "url": meta.get("sourceURL") or url, "pages": 1}
        log(f"[firecrawl] scrape HTTP {r.status_code}")
    except Exception as e:  # noqa: BLE001
        log(f"[firecrawl] scrape error: {e}")
    return None


def _firecrawl_crawl(url, limit, progress=lambda c, t: None, log=lambda m: None, deadline=None):
    """Multi-page crawl via Firecrawl v2 (async: start, then poll). Returns dict or None."""
    key = os.environ.get("FIRECRAWL_API_KEY", "").strip()
    if not key:
        return None
    import requests
    import time
    hdr = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    try:
        r = requests.post("https://api.firecrawl.dev/v2/crawl", headers=hdr,
                          json={"url": url, "limit": limit,
                                "scrapeOptions": {"formats": ["markdown"]}}, timeout=60)
        if r.status_code != 200:
            log(f"[firecrawl] crawl start HTTP {r.status_code}")
            return None
        cid = (r.json() or {}).get("id")
        if not cid:
            return None
    except Exception as e:  # noqa: BLE001
        log(f"[firecrawl] crawl start error: {e}")
        return None
    status_url = f"https://api.firecrawl.dev/v2/crawl/{cid}"
    deadline = deadline or (time.time() + 300)
    pages = []
    while time.time() < deadline:
        try:
            s = requests.get(status_url, headers={"Authorization": f"Bearer {key}"}, timeout=60).json()
        except Exception as e:  # noqa: BLE001
            log(f"[firecrawl] crawl poll error: {e}")
            break
        if s.get("data"):
            pages = s["data"]
        progress(s.get("completed") or len(pages), s.get("total") or limit)
        st = s.get("status")
        if st in ("completed", "failed"):
            break
        time.sleep(3)
    parts = []
    for p in pages:
        md = (p.get("markdown") or "").strip()
        if not md:
            continue
        meta = p.get("metadata") or {}
        parts.append(f"## {meta.get('title') or ''}\n\n<{meta.get('sourceURL') or meta.get('url') or ''}>\n\n{md}")
    if not parts:
        return None
    return {"markdown": "\n\n---\n\n".join(parts), "title": f"Crawl of {url}",
            "url": url, "pages": len(parts)}


def _exa_extract_fallback(url, log=lambda m: None):
    """Last-resort single-page fetch via Exa's cleaned-text contents."""
    try:
        from engines.research.exa_search import exa_contents
        txt = (exa_contents(url, log=log) or "").strip()
        if txt:
            return {"markdown": txt, "title": url, "url": url, "pages": 1}
    except Exception as e:  # noqa: BLE001
        log(f"[exa] extract fallback error: {e}")
    return None


def _extract_worker(job_id, url, mode, limit):
    def setj(**kw):
        with _JOBS_LOCK:
            j = _JOBS.get(job_id)
            if j:
                j.update(kw)

    def log(m):
        setj(message=m)

    try:
        result, provider = None, None
        if mode == "crawl":
            setj(message=f"Crawling {url} (up to {limit} pages) via Firecrawl…", pct=None)
            result = _firecrawl_crawl(
                url, limit, log=log,
                progress=lambda c, t: setj(message=f"Crawling… {c}/{t or limit} pages",
                                           pct=(int(min(c, t) * 100 / t) if t else None)))
            provider = "Firecrawl (crawl)" if result else None
            if not result:
                setj(message="Crawl unavailable — falling back to a single-page scrape…", pct=None)
                result = _firecrawl_scrape(url, log=log)
                provider = "Firecrawl (scrape)" if result else None
        else:
            setj(message=f"Scraping {url} via Firecrawl…", pct=None)
            result = _firecrawl_scrape(url, log=log)
            provider = "Firecrawl" if result else None

        if not result:
            setj(message="Firecrawl failed — trying Exa…", pct=None)
            result = _exa_extract_fallback(url, log=log)
            provider = "Exa (fallback)" if result else None

        if not result:
            setj(done=True, error="Could not fetch this site's contents via Firecrawl or Exa.")
            return
        setj(done=True, pct=100, message="Done.",
             result={"markdown": result["markdown"], "title": result.get("title") or url,
                     "url": result.get("url") or url, "pages": result.get("pages", 1),
                     "provider": provider, "mode": mode})
    except Exception as e:  # noqa: BLE001
        setj(done=True, error=f"{type(e).__name__}: {e}")


@app.route("/api/deep_research/extract", methods=["POST", "OPTIONS"])
def deep_research_extract():
    if request.method == "OPTIONS":
        return "", 204
    url = (request.form.get("url") or "").strip()
    mode = (request.form.get("mode") or "scrape").strip().lower()
    if mode not in ("scrape", "crawl"):
        mode = "scrape"
    try:
        limit = int(request.form.get("limit") or 25)
    except (TypeError, ValueError):
        limit = 25
    limit = max(1, min(limit, 50))          # hard cap on crawl breadth
    if not re.match(r"^https?://", url, re.I):
        return jsonify({"error": "a valid http(s) URL is required"}), 400
    job_id = os.urandom(8).hex()
    with _JOBS_LOCK:
        _JOBS[job_id] = {"stage": "extract", "pct": None, "message": "Starting…",
                         "done": False, "error": None, "result": None}
    threading.Thread(target=_extract_worker, args=(job_id, url, mode, limit), daemon=True).start()
    return jsonify({"job_id": job_id}), 202


@app.route("/api/deep_research/extract_status", methods=["GET"])
def deep_research_extract_status():
    job_id = request.args.get("job", "")
    with _JOBS_LOCK:
        job = _JOBS.get(job_id)
        if not job:
            return jsonify({"error": "unknown job"}), 404
        payload = {"done": job["done"], "error": job["error"],
                   "pct": job["pct"], "message": job["message"]}
        if job["done"] and not job["error"]:
            payload["result"] = job["result"]
        if job["done"]:
            _JOBS.pop(job_id, None)
    return jsonify(payload)


@app.route("/api/deep_research/prompt", methods=["GET", "POST", "OPTIONS"])
def deep_research_prompt():
    """Read/update the governance prompt (prompts/deep_research.md). Loaded fresh at the
    start of every run, so a save applies to all future runs immediately — no restart.
    A rolling .bak is kept; a blank value is ignored so the file can't be wiped.
    GET → {prompt, name}   POST {prompt} → writes it."""
    if request.method == "OPTIONS":
        return "", 204
    path = os.path.join(_PROMPTS_DIR, "deep_research.md")
    if request.method == "GET":
        try:
            with open(path, encoding="utf-8") as fh:
                content = fh.read()
        except Exception:
            content = ""
        return jsonify({"prompt": content, "name": "deep_research.md"})
    data = request.get_json(silent=True) or {}
    val = data.get("prompt")
    if not (isinstance(val, str) and val.strip()):
        return jsonify({"ok": True, "saved": False, "message": "ignored blank prompt"})
    try:
        _backup_then_write(path, val)
    except Exception as e:  # noqa: BLE001
        return jsonify({"error": f"could not save: {type(e).__name__}: {e}"}), 500
    return jsonify({"ok": True, "saved": True})


# ══════════════════════════════════════════════════════════════
#  ODYSSEUS RESEARCH (sub-tool of Deep Research) — vendored comparison engine
#  Alibaba IterResearch-style loop, headless (DuckDuckGo + curl page-fetch) on
#  the same Anthropic key. Here to A/B the methodology vs the visible-Chrome DRT.
# ══════════════════════════════════════════════════════════════

def _ody_msg(ev: dict) -> str:
    ph = ev.get("phase", "")
    rnd = ev.get("round")
    if ph == "planning":
        return "Planning research strategy…"
    if ph == "searching":
        q = ev.get("query_preview")
        return f"Round {rnd}: searching" + (f" — {q}" if q else "…")
    if ph == "reading":
        t = ev.get("title") or ev.get("url")
        rp = f"Round {rnd}: " if rnd else ""
        return f"{rp}reading — {t[:70]}" if t else f"{rp}reading sources…"
    if ph == "analyzing":
        return f"Round {rnd}: synthesizing findings…"
    if ph == "writing":
        return ev.get("message") or "Writing final report…"
    if ph in ("warning", "error"):
        return ev.get("message") or ph
    return ev.get("message") or (ph or "Working…")


def _ody_worker(job_id, query, max_rounds, max_time, category):
    job = _JOBS[job_id]

    def prog(ev):
        msg = _ody_msg(ev)
        with _JOBS_LOCK:
            job["phase"] = ev.get("phase", "")
            job["message"] = msg
            job["stats"] = {"round": ev.get("round"),
                            "sources": ev.get("total_sources"),
                            "findings": ev.get("total_findings")}
        _push_event(job_id, "ody", msg)

    try:
        import asyncio
        from engines.odysseus.deep_research import DeepResearcher
        with _JOBS_LOCK:
            job["message"] = "Starting Odysseus engine…"
        researcher = DeepResearcher(
            llm_endpoint="https://api.anthropic.com/v1/messages",  # ignored by our adapter
            llm_model="claude-opus-4-8",
            max_rounds=max_rounds,
            max_time=max_time,
            progress_callback=prog,
            category=(category or None),
        )
        report = asyncio.run(researcher.research(query))
        stats = researcher.get_stats()
        try:
            docx_b64 = base64.b64encode(
                _memo_to_docx_bytes(report or "", "Odysseus Research")).decode()
        except Exception:
            docx_b64 = ""
        title = _report_title(query, report or "")
        result = {
            "query": query,
            "title": title,
            "report_md": report or "",
            "docx_b64": docx_b64,
            "saved_path": _autosave_report(docx_b64, query, "Odysseus", title=title),
            "stats": stats,
            "sources": researcher.analyzed_urls,
            "source_count": len(researcher.urls_fetched),
        }
        with _JOBS_LOCK:
            job["result"] = result
            job["phase"] = "done"; job["message"] = "Done"; job["done"] = True
    except Exception:
        with _JOBS_LOCK:
            job["error"] = traceback.format_exc()
            job["done"] = True


@app.route("/api/odysseus_research", methods=["POST", "OPTIONS"])
def odysseus_research_start():
    if request.method == "OPTIONS":
        return "", 204
    query = (request.form.get("query") or "").strip()
    if not query:
        return jsonify({"error": "query is required"}), 400
    try:
        max_rounds = int(request.form.get("max_rounds") or 4)
    except (ValueError, TypeError):
        max_rounds = 4
    try:
        max_time = int(request.form.get("max_time") or 300)
    except (ValueError, TypeError):
        max_time = 300
    max_rounds = min(12, max(1, max_rounds))
    max_time = min(1200, max(60, max_time))
    category = (request.form.get("category") or "").strip() or None

    job_id = os.urandom(8).hex()
    with _JOBS_LOCK:
        _JOBS[job_id] = {
            "phase": "planning", "message": "Starting…",
            "stats": {}, "done": False, "error": None, "result": None,
            "events": [], "eseq": 0,
        }
    threading.Thread(target=_ody_worker,
                     args=(job_id, query, max_rounds, max_time, category),
                     daemon=True).start()
    return jsonify({"job_id": job_id}), 202


@app.route("/api/odysseus_research/status", methods=["GET"])
def odysseus_research_status():
    job_id = request.args.get("job", "")
    try:
        after = int(request.args.get("after", 0) or 0)
    except (TypeError, ValueError):
        after = 0
    with _JOBS_LOCK:
        job = _JOBS.get(job_id)
        if not job:
            return jsonify({"error": "unknown job"}), 404
        payload = {
            "phase": job.get("phase"), "message": job.get("message"),
            "stats": job.get("stats", {}), "done": job["done"], "error": job["error"],
            "events": [e for e in (job.get("events") or []) if e[0] > after],
            "eseq": job.get("eseq", 0),
        }
        if job["done"] and not job["error"]:
            payload["result"] = job["result"]
        if job["done"]:
            _JOBS.pop(job_id, None)
    return jsonify(payload)


def _port_already_serving(port: int) -> bool:
    import socket
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.settimeout(0.4)
        return s.connect_ex(("127.0.0.1", port)) == 0


if __name__ == "__main__":
    host = os.environ.get("HOST", "0.0.0.0")
    if _port_already_serving(PORT):
        print(f"dr_server: already serving on port {PORT} — refusing to start a duplicate.")
        sys.exit(0)
    print(f"Deep Research server running on http://{host}:{PORT}")
    app.run(host=host, port=PORT, debug=False, threaded=True)
